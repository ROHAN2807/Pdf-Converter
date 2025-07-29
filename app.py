import os
from flask import Flask, request, render_template, send_file
from pdf2docx import Converter
import fitz
from pptx import Presentation
from pptx.util import Inches
import subprocess
from zipfile import ZipFile
from werkzeug.utils import secure_filename




app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
CONVERTED_FOLDER = "converted"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/convert')
def convert_page():
    option = request.args.get('option')
    return render_template('convert_ui.html', option=option)

@app.route('/convert_pdf', methods=['POST'])
def convert_pdf():
    file = request.files['pdf']
    option = request.form['option']
    filename = file.filename
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(input_path)
    base_name = os.path.splitext(filename)[0]

    if option == 'word':
        output_path = os.path.join(CONVERTED_FOLDER, base_name + '.docx')
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
    elif option == 'images':
        output_folder = os.path.join(CONVERTED_FOLDER, base_name + '_images')
        os.makedirs(output_folder, exist_ok=True)
        doc = fitz.open(input_path)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            pix.save(os.path.join(output_folder, f'page_{i+1}.png'))
        doc.close()
        output_path = output_folder
    elif option == 'ppt':
        output_folder = os.path.join(CONVERTED_FOLDER, base_name + '_slides')
        os.makedirs(output_folder, exist_ok=True)
        doc = fitz.open(input_path)
        image_paths = []
        for i, page in enumerate(doc):
            image_path = os.path.join(output_folder, f'page_{i+1}.png')
            pix = page.get_pixmap(dpi=150)
            pix.save(image_path)
            image_paths.append(image_path)
        doc.close()
        ppt_path = os.path.join(CONVERTED_FOLDER, base_name + '.pptx')
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
        prs.save(ppt_path)
        output_path = ppt_path

    if os.path.isdir(output_path):
        zip_name = os.path.join(CONVERTED_FOLDER, base_name + '_images.zip')
        with ZipFile(zip_name, 'w') as zipf:
            for img in os.listdir(output_path):
                img_path = os.path.join(output_path, img)
                zipf.write(img_path, arcname=img)
        output_path = zip_name

    return send_file(output_path, as_attachment=True)

@app.route('/compress')
def compress_ui():
    return '''
        <form action="/compress_pdf" method="POST" enctype="multipart/form-data">
            <input type="file" name="pdf" required><br><br>
            <button type="submit">Compress</button>
        </form>
    '''

@app.route('/compress_pdf', methods=['POST'])
def compress_pdf():
    file = request.files['pdf']
    filename = secure_filename(file.filename)
    input_path = os.path.join(UPLOAD_FOLDER, filename)
    output_path = os.path.join(CONVERTED_FOLDER, 'compressed_' + filename)
    file.save(input_path)
    cmd = ['gswin64c', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
           '-dPDFSETTINGS=/screen', '-dNOPAUSE', '-dQUIET', '-dBATCH',
           f'-sOutputFile={output_path}', input_path]
    subprocess.run(cmd)
    return send_file(output_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)

